import os
import uuid
import time
import logging
import json
from datetime import datetime
import requests
from flask import Flask, request, jsonify, send_from_directory, render_template, flash, session, redirect, url_for
from flask_cors import CORS
from pptx import Presentation
import fitz  # PyMuPDF for PDF text extraction
import sqlite3
import bcrypt
import tiktoken
import concurrent.futures
from functools import wraps

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger('tutorx')
handler = logging.FileHandler('tutorx.log')
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

app = Flask(__name__, static_folder='static')
app.secret_key = os.getenv('SECRET_KEY', 'supersecretkey')  # For session management
CORS(app)

# Create necessary directories
STATIC_DIR = 'static'
ASSETS_DIR = os.path.join(STATIC_DIR, 'assets')
TEMP_DIR = 'temp'
for directory in [STATIC_DIR, ASSETS_DIR, TEMP_DIR]:
    os.makedirs(directory, exist_ok=True)

# Initialize SQLite database
DB_PATH = 'users.db'

def get_db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            tokens INTEGER DEFAULT 5,
            last_token_reset TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            is_admin INTEGER DEFAULT 0
        )
    ''')
    conn.commit()
    conn.close()

# User management functions
def get_user_by_username(username):
    """Get user by username"""
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE username = ?', (username,)).fetchone()
    conn.close()
    return user

def get_user_by_id(user_id):
    """Get user by ID"""
    conn = get_db_connection()
    user = conn.execute('SELECT * FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    return user

def is_admin(user_id):
    """Check if user has admin privileges"""
    conn = get_db_connection()
    user = conn.execute('SELECT is_admin FROM users WHERE id = ?', (user_id,)).fetchone()
    conn.close()
    return user and user['is_admin'] == 1

def init_ownerx():
    """Initialize OwnerX admin account if it doesn't exist"""
    if not get_user_by_username('OwnerX'):
        password_hash = bcrypt.hashpw(b'adminpass', bcrypt.gensalt())  # Default password: adminpass
        conn = get_db_connection()
        conn.execute(
            'INSERT INTO users (username, password_hash, is_admin) VALUES (?, ?, 1)',
            ('OwnerX', password_hash)
        )
        conn.commit()
        conn.close()
        logger.info('OwnerX admin account created')

# Initialize database and admin account
init_db()
init_ownerx()

# Global processing status
processing_status = {
    'status': None,
    'progress': 0,
    'current_slide': 0,
    'total_slides': 0,
    'current_chunk': 0,
    'total_chunks': 0,
    'flashcard_count': 0,
    'error': None
}

# OpenRouter API configuration
API_CONFIG = {
    'url': 'https://openrouter.ai/api/v1/chat/completions',
    'key': os.getenv('apikey') or os.getenv('APIKEY') or os.getenv('API_KEY') or '',  # Support multiple env var names for API key
    'model': 'mistralai/mistral-7b-instruct:free'  # Changed to a more reliable model
}

def reset_tokens_if_needed(user):
    """Reset tokens to 5 if a week has passed since last reset"""
    conn = get_db_connection()
    cursor = conn.cursor()
    last_reset = user['last_token_reset']
    now = datetime.utcnow()
    if last_reset is None or (now - datetime.strptime(last_reset, '%Y-%m-%d %H:%M:%S')).days >= 7:
        cursor.execute('UPDATE users SET tokens = 5, last_token_reset = ? WHERE id = ?', (now.strftime('%Y-%m-%d %H:%M:%S'), user['id']))
        conn.commit()
    conn.close()

# Decorators
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def admin_required(f):
    """Decorator to require admin access for routes"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session or not is_admin(session['user_id']):
            flash('Admin access required.')
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# Routes
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username'].strip()
        password = request.form['password']
        if not username or not password:
            flash('Username and password are required.')
            return render_template('register.html')
        if get_user_by_username(username):
            flash('Username already exists.')
            return render_template('register.html')
        password_hash = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
        conn = get_db_connection()
        conn.execute('INSERT INTO users (username, password_hash) VALUES (?, ?)', (username, password_hash))
        conn.commit()
        conn.close()
        flash('Registration successful. Please log in.')
        return redirect(url_for('login'))
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username'].strip()
        password = request.form['password']
        user = get_user_by_username(username)
        if user and bcrypt.checkpw(password.encode('utf-8'), user['password_hash']):
            reset_tokens_if_needed(user)
            session['user_id'] = user['id']
            session['username'] = user['username']
            if user['is_admin']:
                return redirect(url_for('admin_dashboard'))
            return redirect(url_for('autox'))
        else:
            flash('Invalid username or password.')
            return render_template('login.html')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/autox.html')
@login_required
def autox():
    user = get_user_by_id(session['user_id'])
    reset_tokens_if_needed(user)
    # Render autox.html template with variables
    return render_template('autox.html', username=user['username'], tokens=user['tokens'])

@app.route('/admin')
@login_required
@admin_required
def admin_dashboard():
    """Admin dashboard to manage user tokens"""
    conn = get_db_connection()
    users = conn.execute('SELECT * FROM users WHERE username != "OwnerX" ORDER BY username').fetchall()
    conn.close()
    return render_template('admin.html', users=users)

@app.route('/admin/add_tokens', methods=['POST'])
@login_required
@admin_required
def add_tokens():
    """Add tokens to a user's account"""
    user_id = request.form.get('user_id')
    tokens = request.form.get('tokens', type=int)
    
    if not user_id or not tokens:
        flash('Invalid request.')
        return redirect(url_for('admin_dashboard'))
    
    conn = get_db_connection()
    conn.execute('UPDATE users SET tokens = tokens + ? WHERE id = ?', (tokens, user_id))
    conn.commit()
    conn.close()
    
    flash(f'Added {tokens} tokens successfully.')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/create_admin', methods=['POST'])
@login_required
@admin_required
def create_admin():
    """Create a new admin user"""
    username = request.form.get('username')
    password = request.form.get('password')
    
    if not username or not password:
        flash('Username and password are required.')
        return redirect(url_for('admin_dashboard'))
    
    if get_user_by_username(username):
        flash('Username already exists.')
        return redirect(url_for('admin_dashboard'))
    
    password_hash = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
    conn = get_db_connection()
    conn.execute(
        'INSERT INTO users (username, password_hash, is_admin) VALUES (?, ?, 1)',
        (username, password_hash)
    )
    conn.commit()
    conn.close()
    
    flash('Admin user created successfully.')
    return redirect(url_for('admin_dashboard'))

@app.route('/process', methods=['POST'])
@login_required
def process_powerpoint():
    user = get_user_by_id(session['user_id'])
    reset_tokens_if_needed(user)
    if user['tokens'] <= 0:
        return jsonify({'error': 'Insufficient tokens', 'message': 'You have no tokens left. Please contact OwnerX to add tokens.'}), 403

    # Deduct 1 token for processing
    conn = get_db_connection()
    conn.execute('UPDATE users SET tokens = tokens - 1 WHERE id = ?', (user['id'],))
    conn.commit()
    conn.close()

    if 'file' not in request.files:
        logger.error("No file provided in request")
        return jsonify({
            'error': 'No file provided',
            'message': 'Please upload a PowerPoint or PDF file'
        }), 400

    file = request.files['file']
    if not file.filename.endswith(('.ppt', '.pptx', '.pdf')):
        logger.error(f"Invalid file format: {file.filename}")
        return jsonify({
            'error': 'Invalid file format',
            'message': 'Please upload a .ppt, .pptx, or .pdf file'
        }), 400

    temp_path = os.path.join(TEMP_DIR, f"{uuid.uuid4()}_{file.filename}")
    logger.info(f"Saving uploaded file to {temp_path}")
    file.save(temp_path)

    try:
        if file.filename.endswith(('.ppt', '.pptx')):
            logger.info("Starting PowerPoint text extraction")
            text_content, slide_count = extract_text_from_pptx(temp_path)
            if not text_content:
                logger.error("No text content extracted from PowerPoint")
                return jsonify({
                    'error': 'Extraction failed',
                    'message': 'Could not extract text from the PowerPoint file'
                }), 400
        elif file.filename.endswith('.pdf'):
            logger.info("Starting PDF text extraction")
            text_content, slide_count = extract_text_from_pdf(temp_path)
            if not text_content:
                logger.error("No text content extracted from PDF")
                return jsonify({
                    'error': 'Extraction failed',
                    'message': 'Could not extract text from the PDF file'
                }), 400
        else:
            logger.error(f"Unsupported file format: {file.filename}")
            return jsonify({
                'error': 'Unsupported file format',
                'message': 'Please upload a supported file type'
            }), 400

        special_instructions = request.form.get('specialInstructions', None)

        logger.info("Starting flashcard generation")
        flashcards = generate_flashcards(text_content, slide_count, special_instructions)
        if not flashcards:
            logger.error("Flashcard generation failed")
            return jsonify({
                'error': 'Generation failed',
                'message': processing_status.get('error', 'Failed to generate flashcards from the content')
            }), 500

        logger.info("Saving generated flashcards")
        flashcard_id = save_flashcards_txt(flashcards, file.filename)
        if not flashcard_id:
            logger.error("Failed to save flashcards")
            return jsonify({
                'error': 'Save failed',
                'message': 'Failed to save generated flashcards'
            }), 500

        return jsonify({
            'success': True,
            'flashcard_id': flashcard_id,
            'count': len(flashcards),
            'flashcards': flashcards,
            'share_link': f'/flashcards.html?id={flashcard_id}'
        })

    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        return jsonify({
            'error': 'Processing failed',
            'message': str(e)
        }), 500

    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)
            logger.info(f"Cleaned up temporary file: {temp_path}")

def update_status(status=None, progress=None, current_slide=None, total_slides=None, 
                 current_chunk=None, total_chunks=None, flashcard_count=None, error=None):
    """Update the global processing status"""
    global processing_status
    if status is not None:
        processing_status['status'] = status
    if progress is not None:
        processing_status['progress'] = progress
    if current_slide is not None:
        processing_status['current_slide'] = current_slide
    if total_slides is not None:
        processing_status['total_slides'] = total_slides
    if current_chunk is not None:
        processing_status['current_chunk'] = current_chunk
    if total_chunks is not None:
        processing_status['total_chunks'] = total_chunks
    if flashcard_count is not None:
        processing_status['flashcard_count'] = flashcard_count
    if error is not None:
        processing_status['error'] = error
    logger.info(f"Status updated: {processing_status}")

def get_api_headers():
    """Get headers for OpenRouter API requests"""
    return {
        'Authorization': f'Bearer {API_CONFIG["key"]}',
        'Content-Type': 'application/json',
        'HTTP-Referer': 'http://localhost:8000',  # Default referer
        'X-Title': 'TutorX'
    }

def chunk_text(text, max_tokens=1500):
    """Split text into chunks based on token count using tiktoken tokenizer."""
    encoding = tiktoken.get_encoding("cl100k_base")
    tokens = encoding.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + max_tokens, len(tokens))
        chunk_tokens = tokens[start:end]
        chunk_text = encoding.decode(chunk_tokens)
        chunks.append(chunk_text)
        start = end
    return chunks

def generate_flashcards(text, slide_count, special_instructions=None):
    """Generate flashcards using OpenRouter API with multiple requests for longer content in parallel"""
    try:
        chunks = chunk_text(text)
        all_flashcards = []
        cards_per_chunk = max(3, slide_count // len(chunks))  # Removed upper cap on cards per chunk
        
        update_status(
            status="Starting flashcard generation",
            progress=0,
            current_chunk=0,
            total_chunks=len(chunks),
            flashcard_count=0
        )

        def process_chunk(i, chunk):
            update_status(
                status=f"Processing chunk {i}/{len(chunks)}",
                progress=(i-1) * 100 // len(chunks),
                current_chunk=i
            )
            
            prompt = f"""Create {cards_per_chunk} educational flashcards from this content. Each flashcard should have a question and answer separated by a colon.

Content:
{chunk}

Format: One flashcard per line, with question and answer separated by a colon.
Example:
What is the powerhouse of the cell?:Mitochondria
What is photosynthesis?:The process by which plants convert sunlight into energy
"""
            if special_instructions:
                prompt += f"\nSpecial Instructions: {special_instructions}\n"

            prompt += f"\nCreate {cards_per_chunk} flashcards now:"

            payload = {
                'model': API_CONFIG['model'],
                'messages': [
                    {
                        'role': 'system',
                        'content': 'You are a teacher creating educational flashcards. Format: question:answer'
                    },
                    {'role': 'user', 'content': prompt}
                ],
                'temperature': 0.7,
                'max_tokens': 1000
            }

            max_retries = 3
            retry_delay = 2
            
            for attempt in range(max_retries):
                try:
                    logger.info(f"Processing chunk {i}/{len(chunks)} (attempt {attempt + 1}/{max_retries})")
                    response = requests.post(
                        API_CONFIG['url'],
                        headers=get_api_headers(),
                        json=payload,
                        timeout=30
                    )
                    
                    if response.status_code == 200:
                        result = response.json()
                        # Log the full API response for debugging
                        logger.info(f"Full API response: {json.dumps(result)}")
                        if result.get('choices'):
                            content = result['choices'][0]['message']['content'].strip()
                            chunk_cards = []
                            
                            # Log the raw API response content for debugging
                            logger.info(f"API Response content: {content}")
                            
                            for line in content.split('\n'):
                                line = line.strip()
                                if ':' in line and not line.startswith('```'):
                                    # Additional validation
                                    parts = line.split(':', 1)
                                    if len(parts) == 2 and all(part.strip() for part in parts):
                                        chunk_cards.append(line)
                            
                            if chunk_cards:
                                return chunk_cards
                            else:
                                logger.error(f"No valid flashcards found in chunk {i} response")
                                if attempt < max_retries - 1:
                                    continue
                        else:
                            logger.error(f"No choices in API response for chunk {i}")
                            if attempt < max_retries - 1:
                                continue
                    else:
                        logger.error(f"API Error for chunk {i}: {response.status_code} - {response.text}")
                        if attempt < max_retries - 1:
                            update_status(status=f"Retrying chunk {i} (attempt {attempt + 2}/{max_retries})")
                            time.sleep(retry_delay)
                            retry_delay *= 2
                            continue
                        
                except requests.exceptions.RequestException as e:
                    logger.error(f"Request failed for chunk {i}: {str(e)}")
                    if attempt < max_retries - 1:
                        update_status(status=f"Retrying chunk {i} (attempt {attempt + 2}/{max_retries})")
                        time.sleep(retry_delay)
                        retry_delay *= 2
                        continue
                    raise
            return []

        with concurrent.futures.ThreadPoolExecutor() as executor:
            futures = [executor.submit(process_chunk, i, chunk) for i, chunk in enumerate(chunks, 1)]
            for i, future in enumerate(concurrent.futures.as_completed(futures), 1):
                chunk_cards = future.result()
                if chunk_cards:
                    all_flashcards.extend(chunk_cards)
                    update_status(
                        flashcard_count=len(all_flashcards),
                        progress=i * 100 // len(chunks)
                    )

        if not all_flashcards:
            raise Exception("The server is busy, try again in ~15m")

        update_status(
            status="Flashcard generation complete",
            progress=100,
            flashcard_count=len(all_flashcards)
        )
        
        logger.info(f"Successfully generated {len(all_flashcards)} total flashcards")
        return all_flashcards

    except Exception as e:
        logger.error(f"Unexpected error in generate_flashcards: {str(e)}")
        update_status(error=str(e))
        return None

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint slides"""
    try:
        logger.info(f"Opening PowerPoint file: {file_path}")
        prs = Presentation(file_path)
        text_content = []
        slide_count = len(prs.slides)
        
        update_status(
            status="Starting PowerPoint extraction",
            total_slides=slide_count,
            current_slide=0
        )
        
        for i, slide in enumerate(prs.slides, 1):
            update_status(
                status=f"Processing slide {i}/{slide_count}",
                current_slide=i
            )
            
            slide_text = []
            logger.info(f"Processing slide {i}/{slide_count}")
            
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
                slide_text.append(f"Title: {title_text}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text not in slide_text:
                        slide_text.append(text)
            
            if slide_text:
                text_content.append("\n".join(slide_text))
        
        full_text = "\n\n".join(text_content)
        logger.info(f"Total extracted text length: {len(full_text)} characters")
        
        if not full_text.strip():
            raise Exception("No text content found in the PowerPoint file")
            
        return full_text, slide_count
        
    except Exception as e:
        logger.error(f"PowerPoint extraction error: {str(e)}")
        update_status(error=str(e))
        return None, 0

def extract_text_from_pdf(file_path):
    """Extract text from PDF file with improved formatting and structure"""
    try:
        logger.info(f"Opening PDF file: {file_path}")
        doc = fitz.open(file_path)
        text_content = []
        slide_count = doc.page_count
        
        update_status(
            status="Starting PDF extraction",
            total_slides=slide_count,
            current_slide=0
        )
        
        for i, page in enumerate(doc, 1):
            update_status(
                status=f"Processing page {i}/{slide_count}",
                current_slide=i
            )
            
            # Get text blocks with their formatting
            blocks = page.get_text("dict")["blocks"]
            page_content = []
            
            # Process each text block
            for block in blocks:
                if block.get("type") == 0:  # Type 0 is text
                    for line in block.get("lines", []):
                        line_text = ""
                        for span in line.get("spans", []):
                            text = span.get("text", "").strip()
                            if text:
                                # Check if it's a heading based on font size
                                if span.get("size", 0) > 12:  # Adjust this threshold as needed
                                    text = f"\nHeading: {text}\n"
                                line_text += text + " "
                        if line_text.strip():
                            page_content.append(line_text.strip())
            
            # Join the page content with proper spacing
            if page_content:
                page_text = "\n".join(page_content)
                # Add page number for better context
                text_content.append(f"Page {i}:\n{page_text}")
        
        # Join all pages with clear separation
        full_text = "\n\n" + "="*40 + "\n\n".join(text_content) + "\n" + "="*40 + "\n"
        logger.info(f"Total extracted text length: {len(full_text)} characters")
        
        if not full_text.strip():
            raise Exception("No text content found in the PDF file")
        
        # Clean up the text
        full_text = (full_text
            .replace('\u2019', "'")  # Smart quotes
            .replace('\u2018', "'")
            .replace('\u201c', '"')
            .replace('\u201d', '"')
            .replace('\u2013', '-')  # En dash
            .replace('\u2014', '--')  # Em dash
        )
            
        return full_text, slide_count
        
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        update_status(error=str(e))
        return None, 0

def save_flashcards_txt(flashcards, filename):
    """Save flashcards to a text file"""
    try:
        update_status(status="Saving flashcards")
        flashcard_id = str(uuid.uuid4())
        output_path = os.path.join(ASSETS_DIR, f'flashcards_{flashcard_id}.txt')
        
        with open(output_path, 'w') as f:
            f.write('\n'.join(flashcards))
            
        logger.info(f"Saved {len(flashcards)} flashcards to {output_path}")
        return flashcard_id
        
    except Exception as e:
        logger.error(f"Error saving flashcards: {str(e)}")
        update_status(error=str(e))
        return None

@app.route('/')
def index():
    """Serve the landing page"""
    return send_from_directory('.', 'index.html')

@app.route('/flashcards.html')
def flashcards():
    """Serve the flashcards page"""
    return send_from_directory('.', 'flashcards.html')

@app.route('/process/status')
def get_status():
    """Get the current processing status"""
    return jsonify(processing_status)

@app.route('/flashcards/<flashcard_id>', methods=['GET'])
def get_flashcards(flashcard_id):
    """API endpoint to fetch flashcards JSON by flashcard_id"""
    try:
        flashcard_file = os.path.join(ASSETS_DIR, f'flashcards_{flashcard_id}.txt')
        if not os.path.isfile(flashcard_file):
            return jsonify({'error': 'Flashcards not found'}), 404
        
        with open(flashcard_file, 'r', encoding='utf-8') as f:
            lines = f.read().splitlines()
        
        flashcards = []
        for line in lines:
            if ':' in line:
                question, answer = line.split(':', 1)
                flashcards.append({'question': question.strip(), 'answer': answer.strip()})
        
        return jsonify({'flashcards': flashcards})
    except Exception as e:
        logger.error(f"Error fetching flashcards {flashcard_id}: {str(e)}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/static/<path:filename>')
def serve_static(filename):
    """Serve static files"""
    return send_from_directory('static', filename)

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=7080)
